from __future__ import annotations

import functools
import logging
from pathlib import Path
from typing import Any

from semanticfs.chunker import FileChunk, chunk_file_content

logger = logging.getLogger(__name__)

class Embedder:
    """Embeds text and files into vectors, with support for local fine-tuned models, rich PPTX/XLSX extraction, and dynamic semantic chunking."""
    def __init__(self, model_name: str = "all-MiniLM-L6-v2", max_tokens: int = 512):
        self.model_name = model_name
        self.max_tokens = max_tokens
        self._model: Any = None

        # Check for custom fine-tuned model
        custom_model_dir = Path("~/.semanticfs/custom_model").expanduser()
        if custom_model_dir.exists():
            self.model_path = str(custom_model_dir)
        else:
            self.model_path = model_name

    @property
    def model(self):
        if self._model is None:
            logger.info(f"Loading embedding model from '{self.model_path}'...")
            from sentence_transformers import SentenceTransformer
            self._model = SentenceTransformer(self.model_path)
        return self._model

    @functools.lru_cache(maxsize=512)
    def embed_text(self, text: str) -> list[float]:
        """Embed a single text string with LRU caching."""
        truncated = text[:self.max_tokens * 4] 
        embedding = self.model.encode(truncated)
        return embedding.tolist()

    def embed_batch(self, texts: list[str]) -> list[list[float]]:
        """Batch embed multiple texts."""
        truncated_texts = [text[:self.max_tokens * 4] for text in texts]
        embeddings = self.model.encode(truncated_texts)
        return embeddings.tolist()

    def extract_chunks(self, filepath: Path) -> list[FileChunk]:
        """Extract content and dynamically split file into semantic chunks with line tracking."""
        content = self._extract_content(filepath)
        return chunk_file_content(filepath, content)

    def _extract_content(self, filepath: Path) -> str:
        """Extract text from text, code, document, slides, tabular, or binary files gracefully."""
        ext = filepath.suffix.lower()
        content = ""
        try:
            if ext == ".pdf":
                import fitz  # PyMuPDF
                doc = fitz.open(filepath)
                for page in doc:
                    content += page.get_text() + "\n"
                doc.close()
            elif ext == ".docx":
                import docx
                doc = docx.Document(filepath)
                content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(filepath)
                slide_texts = []
                for i, slide in enumerate(prs.slides, start=1):
                    lines = [f"--- Slide {i} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            lines.append(shape.text.strip())
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            lines.append(f"Notes: {notes}")
                    slide_texts.append("\n".join(lines))
                content = "\n\n".join(slide_texts)
            elif ext in (".xlsx", ".xls"):
                import openpyxl
                wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
                sheet_texts = []
                for sheetname in wb.sheetnames:
                    ws = wb[sheetname]
                    lines = [f"--- Sheet: {sheetname} ---"]
                    for row in ws.iter_rows(values_only=True):
                        row_vals = [str(val) for val in row if val is not None]
                        if row_vals:
                            lines.append(" | ".join(row_vals))
                    sheet_texts.append("\n".join(lines[:100]))  # Cap at 100 rows per sheet for speed
                content = "\n\n".join(sheet_texts)
            elif ext in (".csv", ".tsv"):
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read(5000)
            elif ext in (".ipynb", ".json", ".yaml", ".yml", ".toml", ".xml"):
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read(5000)
            elif ext in (".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".c", ".cpp", ".h", ".hpp", ".cs", ".java", ".rs", ".go", ".rb", ".php", ".swift", ".kt", ".md", ".txt", ".sh", ".bat", ".ps1", ".sql", ".r", ".R", ".log"):
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read(10000)
            else:
                stat = filepath.stat() if filepath.exists() else None
                size_kb = round(stat.st_size / 1024, 1) if stat else 0
                content = f"Media/Binary File: {filepath.name}, Type: {ext}, Folder: {filepath.parent.name}, Size: {size_kb} KB"
        except Exception as e:
            logger.debug(f"Error extracting content from {filepath}: {e}")
            content = f"File: {filepath.name}, Type: {ext}, Folder: {filepath.parent.name}"
        return content
